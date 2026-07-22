"""Text extraction (F06) + OCR fallback (F07).

PDF  -> Azure AI Document Intelligence (prebuilt-read) รองรับทั้ง text + scanned (OCR)
PPTX -> python-pptx (text-based); Document Intelligence ไม่รองรับ .pptx ตรง ๆ
"""
from __future__ import annotations

import io
import os

from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential


def _docintel_client() -> DocumentAnalysisClient:
    return DocumentAnalysisClient(
        endpoint=os.environ["DOCINTEL_ENDPOINT"],
        credential=AzureKeyCredential(os.environ["DOCINTEL_KEY"]),
    )


def extract_pdf(data: bytes) -> str:
    """PDF -> text. prebuilt-read จัดการ OCR ให้อัตโนมัติถ้าเป็น scanned (F07)."""
    client = _docintel_client()
    poller = client.begin_analyze_document("prebuilt-read", document=data)
    result = poller.result()
    return "\n".join(line.content for page in result.pages for line in page.lines)


def extract_pptx(data: bytes) -> str:
    """PPTX -> text ต่อ slide (คงลำดับ slide ไว้เพื่อ map เข้า skeleton)."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = [
            shape.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        if parts:
            chunks.append(f"[Slide {i}]\n" + "\n".join(parts))
    return "\n\n".join(chunks)


def extract_text(data: bytes, content_type: str, filename: str) -> str:
    """Dispatch ตาม type. Raise ถ้า format ไม่รองรับ (F03 validation)."""
    name = filename.lower()
    if content_type == "application/pdf" or name.endswith(".pdf"):
        return extract_pdf(data)
    if name.endswith(".pptx") or "presentation" in content_type:
        return extract_pptx(data)
    raise ValueError(f"Unsupported format: {content_type} / {filename}")
